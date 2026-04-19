"""Render PPTPlan into a local .pptx file."""

from datetime import datetime
from pathlib import Path

from pptx import Presentation

from schemas.ppt_plan_schema import PPTPlan


def _safe_name(name: str) -> str:
	cleaned = "".join(ch for ch in name if ch.isalnum() or ch in ("-", "_"))
	return cleaned[:40] or "im2ppt"


def render_plan_to_pptx(plan: PPTPlan, output_dir: str = "temp") -> str:
	"""Render a minimal presentation and return absolute file path."""
	out_dir = Path(output_dir)
	out_dir.mkdir(parents=True, exist_ok=True)

	ts = datetime.now().strftime("%Y%m%d_%H%M%S")
	file_name = f"{_safe_name(plan.title)}_{ts}.pptx"
	file_path = out_dir / file_name

	prs = Presentation()

	# Cover slide
	cover_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]
	cover = prs.slides.add_slide(cover_layout)
	if cover.shapes.title:
		cover.shapes.title.text = plan.title
	if len(cover.placeholders) > 1:
		cover.placeholders[1].text = f"IM2PPT 自动生成 | {datetime.now().strftime('%Y-%m-%d %H:%M')}"

	# Content slides
	content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
	for slide_plan in plan.slides:
		slide = prs.slides.add_slide(content_layout)
		if slide.shapes.title:
			slide.shapes.title.text = slide_plan.title

		if len(slide.placeholders) > 1:
			tf = slide.placeholders[1].text_frame
			tf.clear()
			if slide_plan.bullets:
				tf.text = slide_plan.bullets[0]
				for bullet in slide_plan.bullets[1:]:
					p = tf.add_paragraph()
					p.text = bullet
					p.level = 0
			else:
				tf.text = "待补充"

	prs.save(str(file_path))
	return str(file_path.resolve())
